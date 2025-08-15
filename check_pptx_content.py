#!/usr/bin/env python3
"""
Script to check the content of generated PowerPoint files
"""

from pptx import Presentation
import os

def check_pptx_content(pptx_path: str):
    """Check the content of a PowerPoint file"""
    print(f"📊 Checking PowerPoint content: {pptx_path}")
    print("=" * 60)
    
    if not os.path.exists(pptx_path):
        print(f"❌ PowerPoint file not found: {pptx_path}")
        return
    
    try:
        prs = Presentation(pptx_path)
        print(f"✅ PowerPoint opened successfully")
        print(f"   Total slides: {len(prs.slides)}")
        print()
        
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"📄 Slide {slide_num}:")
            
            # Check title
            if slide.shapes.title:
                title = slide.shapes.title.text
                if title:
                    print(f"   Title: {title}")
            
            # Check content
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text and text != title:
                            print(f"   Content: {text[:100]}{'...' if len(text) > 100 else ''}")
            
            print()
        
        # File info
        file_size = os.path.getsize(pptx_path)
        print(f"📊 File Information:")
        print(f"   Size: {file_size} bytes ({file_size/1024:.1f} KB)")
        print(f"   Slides: {len(prs.slides)}")
        
    except Exception as e:
        print(f"❌ Error reading PowerPoint: {e}")

def main():
    """Check all generated PowerPoint files"""
    print("PowerPoint Content Checker")
    print("=" * 60)
    
    # Check the user's converted file
    user_pptx = "outputs/pptx/user_document_conversion.pptx"
    if os.path.exists(user_pptx):
        check_pptx_content(user_pptx)
        print("\n" + "=" * 60)
    
    # Check the sample converted file
    sample_pptx = "outputs/pptx/sample_conversion.pptx"
    if os.path.exists(sample_pptx):
        check_pptx_content(sample_pptx)
        print("\n" + "=" * 60)
    
    # Check the test converted file
    test_pptx = "outputs/pptx/test_conversion.pptx"
    if os.path.exists(test_pptx):
        check_pptx_content(test_pptx)

if __name__ == "__main__":
    main()
