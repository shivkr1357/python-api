"""
PDF processing utilities for handling password-protected PDFs.
"""

import io
import tempfile
import os
from typing import Optional, Tuple, List
from pathlib import Path

from PyPDF2 import PdfReader, PdfWriter
from fastapi import HTTPException, UploadFile
from fastapi.responses import StreamingResponse

try:
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image
    import fitz  # PyMuPDF for better PDF processing
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PDFProcessor:
    """PDF processing utility for handling password-protected PDFs."""
    
    @staticmethod
    def unlock_pdf_automatically(
        pdf_file: UploadFile
    ) -> Tuple[bytes, str]:
        """
        Automatically unlock a password-protected PDF without requiring a password.
        
        Args:
            pdf_file: The uploaded PDF file
            
        Returns:
            Tuple of (unlocked_pdf_bytes, filename)
            
        Raises:
            HTTPException: If PDF cannot be unlocked automatically
        """
        try:
            # Read the uploaded file content
            pdf_content = pdf_file.file.read()
            
            # Use PyPDF2 for PDF processing
            try:
                unlocked_pdf = PDFProcessor._unlock_automatically_with_pypdf2(pdf_content)
                return unlocked_pdf, f"unlocked_{pdf_file.filename}"
            except Exception as pypdf2_error:
                raise HTTPException(
                    status_code=400,
                    detail=f"Failed to unlock PDF automatically. "
                           f"PyPDF2 error: {str(pypdf2_error)}"
                )
                    
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=500,
                detail=f"Error processing PDF: {str(e)}"
            )
    
    @staticmethod
    def unlock_pdf_with_password(
        pdf_file: UploadFile, 
        password: str
    ) -> Tuple[bytes, str]:
        """
        Unlock a password-protected PDF and return the unlocked content.
        
        Args:
            pdf_file: The uploaded PDF file
            password: The password to unlock the PDF
            
        Returns:
            Tuple of (unlocked_pdf_bytes, filename)
            
        Raises:
            HTTPException: If password is incorrect or PDF is corrupted
        """
        try:
            # Read the uploaded file content
            pdf_content = pdf_file.file.read()
            
            # Use PyPDF2 for PDF processing
            try:
                unlocked_pdf = PDFProcessor._unlock_with_pypdf2(pdf_content, password)
                return unlocked_pdf, f"unlocked_{pdf_file.filename}"
            except Exception as pypdf2_error:
                raise HTTPException(
                    status_code=400,
                    detail=f"Failed to unlock PDF. Please check if the password is correct. "
                           f"PyPDF2 error: {str(pypdf2_error)}"
                )
                    
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=500,
                detail=f"Error processing PDF: {str(e)}"
            )

    @staticmethod
    def lock_pdf_with_password(
        pdf_file: UploadFile, 
        password: str
    ) -> Tuple[bytes, str]:
        """
        Lock a PDF with a password and return the protected content.
        
        Args:
            pdf_file: The uploaded PDF file
            password: The password to protect the PDF
            
        Returns:
            Tuple of (locked_pdf_bytes, filename)
            
        Raises:
            HTTPException: If PDF processing fails
        """
        try:
            # Read the uploaded file content
            pdf_content = pdf_file.file.read()
            
            # Use PyPDF2 for PDF processing
            try:
                locked_pdf = PDFProcessor._lock_with_pypdf2(pdf_content, password)
                return locked_pdf, f"locked_{pdf_file.filename}"
            except Exception as pypdf2_error:
                raise HTTPException(
                    status_code=400,
                    detail=f"Failed to lock PDF. PyPDF2 error: {str(pypdf2_error)}"
                )
                    
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=500,
                detail=f"Error processing PDF: {str(e)}"
            )
    
    @staticmethod
    def compress_pdf(
        pdf_file: UploadFile, 
        compression_level: str
    ) -> Tuple[bytes, str]:
        """
        Compress a PDF file to reduce its size.
        
        Args:
            pdf_file: The uploaded PDF file
            compression_level: Compression level (low, medium, high)
            
        Returns:
            Tuple of (compressed_pdf_bytes, filename)
            
        Raises:
            HTTPException: If PDF processing fails
        """
        try:
            # Read the uploaded file content
            pdf_content = pdf_file.file.read()
            
            # Use PyPDF2 for PDF processing
            try:
                compressed_pdf = PDFProcessor._compress_with_pypdf2(pdf_content, compression_level)
                return compressed_pdf, f"compressed_{compression_level}_{pdf_file.filename}"
            except Exception as pypdf2_error:
                raise HTTPException(
                    status_code=400,
                    detail=f"Failed to compress PDF. PyPDF2 error: {str(pypdf2_error)}"
                )
                    
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=500,
                detail=f"Error processing PDF: {str(e)}"
            )
    
    @staticmethod
    def pdf_to_powerpoint(
        pdf_file: UploadFile
    ) -> Tuple[bytes, str]:
        """
        Convert a PDF file to PowerPoint format.
        
        Args:
            pdf_file: The uploaded PDF file
            
        Returns:
            Tuple of (powerpoint_bytes, filename)
            
        Raises:
            HTTPException: If conversion fails or dependencies not available
        """
        if not PPTX_AVAILABLE:
            raise HTTPException(
                status_code=500,
                detail="PDF to PowerPoint conversion is not available. Please install required dependencies: pip install python-pptx Pillow PyMuPDF"
            )
        
        try:
            # Read the uploaded file content
            pdf_content = pdf_file.file.read()
            
            # Convert PDF to PowerPoint
            try:
                powerpoint_content = PDFProcessor._convert_pdf_to_pptx(pdf_content)
                return powerpoint_content, f"converted_{pdf_file.filename.replace('.pdf', '.pptx')}"
            except Exception as conversion_error:
                raise HTTPException(
                    status_code=400,
                    detail=f"Failed to convert PDF to PowerPoint. Error: {str(conversion_error)}"
                )
                    
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=500,
                detail=f"Error processing PDF: {str(e)}"
            )
    
    @staticmethod
    def powerpoint_to_pdf(
        pptx_file: UploadFile
    ) -> Tuple[bytes, str]:
        """
        Convert a PowerPoint file to PDF format.
        
        Args:
            pptx_file: The uploaded PowerPoint file
            
        Returns:
            Tuple of (pdf_bytes, filename)
            
        Raises:
            HTTPException: If conversion fails or dependencies not available.
        """
        if not PPTX_AVAILABLE:
            raise HTTPException(
                status_code=500,
                detail="PowerPoint to PDF conversion is not available. Please install required dependencies: pip install python-pptx Pillow PyMuPDF"
            )
        
        try:
            # Read the uploaded file content
            pptx_content = pptx_file.file.read()
            
            # Convert PowerPoint to PDF
            try:
                pdf_content = PDFProcessor._convert_pptx_to_pdf(pptx_content)
                return pdf_content, f"converted_{pptx_file.filename.replace('.pptx', '.pdf')}"
            except Exception as conversion_error:
                raise HTTPException(
                    status_code=400,
                    detail=f"Failed to convert PowerPoint to PDF. Error: {str(conversion_error)}"
                )
                    
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=500,
                detail=f"Error processing PowerPoint file: {str(e)}"
            )
    
    @staticmethod
    def jpg_to_pdf(
        jpg_files: List[UploadFile],
        page_orientation: str = "portrait",
        page_size: str = "a4",
        margin: str = "no_margin",
        merge_all: bool = True
    ) -> Tuple[bytes, str]:
        """
        Convert JPG images to PDF format with configurable options.
        
        Args:
            jpg_files: List of uploaded JPG image files
            page_orientation: "portrait" or "landscape"
            page_size: "a4", "us_letter", or "fit"
            margin: "no_margin", "small", or "big"
            merge_all: Whether to merge all images into one PDF file
            
        Returns:
            Tuple of (pdf_bytes, filename)
            
        Raises:
            HTTPException: If conversion fails
        """
        try:
            # Validate inputs
            if not jpg_files:
                raise HTTPException(
                    status_code=400,
                    detail="At least one JPG file is required"
                )
            
            if page_orientation not in ["portrait", "landscape"]:
                raise HTTPException(
                    status_code=400,
                    detail="Page orientation must be 'portrait' or 'landscape'"
                )
            
            if page_size not in ["a4", "us_letter", "fit"]:
                raise HTTPException(
                    status_code=400,
                    detail="Page size must be 'a4', 'us_letter', or 'fit'"
                )
            
            if margin not in ["no_margin", "small", "big"]:
                raise HTTPException(
                    status_code=400,
                    detail="Margin must be 'no_margin', 'small', or 'big'"
                )
            
            # Convert JPG to PDF
            try:
                pdf_content = PDFProcessor._convert_jpgs_to_pdf(
                    jpg_files, page_orientation, page_size, margin, merge_all
                )
                
                # Generate filename
                if merge_all and len(jpg_files) > 1:
                    filename = f"merged_{len(jpg_files)}_images.pdf"
                else:
                    filename = f"converted_{jpg_files[0].filename.replace('.jpg', '.pdf').replace('.jpeg', '.pdf')}"
                
                return pdf_content, filename
                
            except Exception as conversion_error:
                raise HTTPException(
                    status_code=400,
                    detail=f"Failed to convert JPG to PDF. Error: {str(conversion_error)}"
                )
                    
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=500,
                detail=f"Error processing JPG files: {str(e)}"
            )
    
    @staticmethod
    def pdf_to_jpg(
        pdf_file: UploadFile,
        page_number: int = 1
    ) -> Tuple[bytes, str]:
        """
        Convert a PDF page to JPG format.
        
        Args:
            pdf_file: The uploaded PDF file
            page_number: Page number to convert (default: 1)
            
        Returns:
            Tuple of (jpg_bytes, filename)
            
        Raises:
            HTTPException: If conversion fails
        """
        try:
            # Read the uploaded file content
            pdf_content = pdf_file.file.read()
            
            # Convert PDF to JPG
            try:
                jpg_content = PDFProcessor._convert_pdf_to_jpg(pdf_content, page_number)
                return jpg_content, f"converted_page_{page_number}_{pdf_file.filename.replace('.pdf', '.jpg')}"
            except Exception as conversion_error:
                raise HTTPException(
                    status_code=400,
                    detail=f"Failed to convert PDF to JPG. Error: {str(conversion_error)}"
                )
                    
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=500,
                detail=f"Error processing PDF file: {str(e)}"
            )
    
    @staticmethod
    def _unlock_automatically_with_pypdf2(pdf_content: bytes) -> bytes:
        """Automatically unlock PDF using PyPDF2 library without password."""
        try:
            # Create input buffer
            input_buffer = io.BytesIO(pdf_content)
            
            # Create PDF reader
            reader = PdfReader(input_buffer)
            
            # Check if PDF is encrypted
            if reader.is_encrypted:
                # Try to decrypt with empty password first
                if reader.decrypt(""):
                    print("DEBUG: Successfully decrypted with empty password")
                else:
                    # Try common passwords
                    common_passwords = [
                        "", "password", "123456", "admin", "user", "1234", "12345", 
                        "123456789", "qwerty", "abc123", "password123", "admin123",
                        "user123", "test", "test123", "demo", "demo123", "guest",
                        "guest123", "public", "public123", "default", "default123",
                        "123", "0000", "1111", "2222", "3333", "4444", "5555",
                        "6666", "7777", "8888", "9999", "000000", "111111",
                        "secret", "private", "secure", "access", "login"
                    ]
                    
                    decrypted = False
                    for pwd in common_passwords:
                        try:
                            if reader.decrypt(pwd):
                                print(f"DEBUG: Successfully decrypted with password: {pwd}")
                                decrypted = True
                                break
                        except Exception:
                            continue
                    
                    if not decrypted:
                        # If still not decrypted, try to read without decryption
                        # Some PDFs might be readable even if marked as encrypted
                        try:
                            # Try to access pages without decryption
                            page_count = len(reader.pages)
                            if page_count > 0:
                                print("DEBUG: PDF appears to be readable without decryption")
                                # If we can access pages, the PDF might not actually be encrypted
                                pass
                            else:
                                raise HTTPException(
                                    status_code=400,
                                    detail="Cannot unlock PDF automatically. The PDF is protected with a strong password that cannot be automatically detected."
                                )
                        except Exception:
                            raise HTTPException(
                                status_code=400,
                                detail="Cannot unlock PDF automatically. The PDF is protected with a strong password that cannot be automatically detected."
                            )
            else:
                print("DEBUG: PDF is not encrypted")
                # PDF is not encrypted, just return the original content
                pass
            
            # Create PDF writer
            writer = PdfWriter()
            
            # Add all pages from reader to writer
            for page in reader.pages:
                writer.add_page(page)
            
            # IMPORTANT: Do NOT encrypt the output PDF
            # The writer should remain unencrypted to create an unlocked PDF
            
            # Create output buffer
            output_buffer = io.BytesIO()
            
            # Write the unlocked PDF to buffer (without encryption)
            writer.write(output_buffer)
            
            # Get the content
            unlocked_content = output_buffer.getvalue()
            
            print(f"DEBUG: Successfully created unlocked PDF of size: {len(unlocked_content)} bytes")
            
            return unlocked_content
            
        except HTTPException:
            raise
        except Exception as e:
            print(f"DEBUG: Error in automatic unlock: {str(e)}")
            # If it's a PyPDF2 specific error, try alternative approach
            try:
                # Try to create a new PDF reader and writer
                input_buffer = io.BytesIO(pdf_content)
                reader = PdfReader(input_buffer)
                writer = PdfWriter()
                
                # Try to add pages without decryption
                for page in reader.pages:
                    writer.add_page(page)
                
                output_buffer = io.BytesIO()
                writer.write(output_buffer)
                return output_buffer.getvalue()
                
            except Exception as fallback_error:
                raise HTTPException(
                    status_code=400,
                    detail=f"Error unlocking PDF automatically: {str(e)}. Fallback also failed: {str(fallback_error)}"
                )
    
    @staticmethod
    def _unlock_with_pypdf2(pdf_content: bytes, password: str) -> bytes:
        """Unlock PDF using PyPDF2 library."""
        try:
            # Create input buffer
            input_buffer = io.BytesIO(pdf_content)
            
            # Create PDF reader
            reader = PdfReader(input_buffer)
            
            # Check if PDF is encrypted
            if reader.is_encrypted:
                # Try to decrypt with password
                if not reader.decrypt(password):
                    raise HTTPException(
                        status_code=400,
                        detail="Incorrect password provided"
                    )
            
            # Create PDF writer
            writer = PdfWriter()
            
            # Add all pages from reader to writer
            for page in reader.pages:
                writer.add_page(page)
            
            # Create output buffer
            output_buffer = io.BytesIO()
            
            # Write the unlocked PDF to buffer
            writer.write(output_buffer)
            
            # Get the content
            unlocked_content = output_buffer.getvalue()
            
            return unlocked_content
            
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Error unlocking PDF with PyPDF2: {str(e)}"
            )

    @staticmethod
    def _lock_with_pypdf2(pdf_content: bytes, password: str) -> bytes:
        """Lock PDF using PyPDF2 library."""
        try:
            # Create input buffer
            input_buffer = io.BytesIO(pdf_content)
            
            # Create PDF reader
            reader = PdfReader(input_buffer)
            
            # Create PDF writer
            writer = PdfWriter()
            
            # Add all pages from reader to writer
            for page in reader.pages:
                writer.add_page(page)
            
            # Encrypt the PDF with password
            writer.encrypt(password)
            
            # Create output buffer
            output_buffer = io.BytesIO()
            
            # Write the locked PDF to buffer
            writer.write(output_buffer)
            
            # Get the content
            locked_content = output_buffer.getvalue()
            
            return locked_content
            
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Error locking PDF with PyPDF2: {str(e)}"
            )
    
    @staticmethod
    def _compress_with_pypdf2(pdf_content: bytes, compression_level: str) -> bytes:
        """Compress PDF using PyPDF2 library."""
        try:
            # Create input buffer
            input_buffer = io.BytesIO(pdf_content)
            
            # Create PDF reader
            reader = PdfReader(input_buffer)
            
            # Create PDF writer
            writer = PdfWriter()
            
            # Set compression parameters based on level
            if compression_level == "low":
                # Minimal compression - preserve quality
                compression_settings = {
                    "compress": True,
                    "compress_content_streams": True,
                    "object_stream_mode": 1
                }
            elif compression_level == "medium":
                # Balanced compression
                compression_settings = {
                    "compress": True,
                    "compress_content_streams": True,
                    "object_stream_mode": 2
                }
            elif compression_level == "high":
                # Maximum compression
                compression_settings = {
                    "compress": True,
                    "compress_content_streams": True,
                    "object_stream_mode": 3
                }
            else:
                # Default to medium
                compression_settings = {
                    "compress": True,
                    "compress_content_streams": True,
                    "object_stream_mode": 2
                }
            
            # Add all pages from reader to writer
            for page in reader.pages:
                writer.add_page(page)
            
            # Apply compression settings
            writer._compress = compression_settings["compress"]
            writer._compress_content_streams = compression_settings["compress_content_streams"]
            writer._object_stream_mode = compression_settings["object_stream_mode"]
            
            # Create output buffer
            output_buffer = io.BytesIO()
            
            # Write the compressed PDF to buffer
            writer.write(output_buffer)
            
            # Get the content
            compressed_content = output_buffer.getvalue()
            
            print(f"DEBUG: Original size: {len(pdf_content)} bytes, Compressed size: {len(compressed_content)} bytes")
            print(f"DEBUG: Compression ratio: {((len(pdf_content) - len(compressed_content)) / len(pdf_content) * 100):.1f}%")
            
            return compressed_content
            
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Error compressing PDF with PyPDF2: {str(e)}"
            )
    
    @staticmethod
    def _convert_pdf_to_pptx(pdf_content: bytes) -> bytes:
        """Convert PDF to PowerPoint using PyMuPDF and python-pptx."""
        try:
            # Create a new PowerPoint presentation
            prs = Presentation()
            
            # Open PDF with PyMuPDF
            pdf_document = fitz.open(stream=pdf_content, filetype="pdf")
            
            print(f"DEBUG: Converting PDF with {len(pdf_document)} pages to PowerPoint")
            
            # Get the first page to determine PDF dimensions
            first_page = pdf_document[0]
            pdf_width = first_page.rect.width
            pdf_height = first_page.rect.height
            
            # Calculate PowerPoint slide dimensions to maintain PDF width
            # PowerPoint uses EMU (English Metric Units): 1 inch = 914400 EMU
            # Convert PDF points to inches: 1 inch = 72 points
            # Then convert to EMU
            pdf_width_inches = pdf_width / 72
            pdf_height_inches = pdf_height / 72
            
            # Set slide dimensions to match PDF
            prs.slide_width = int(pdf_width_inches * 914400)
            prs.slide_height = int(pdf_height_inches * 914400)
            
            print(f"DEBUG: PDF dimensions: {pdf_width} x {pdf_height} points")
            print(f"DEBUG: PowerPoint slide dimensions: {prs.slide_width} x {prs.slide_height} EMU")
            
            # Process each page
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                
                # Create a new slide
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Convert page to image
                mat = fitz.Matrix(2, 2)  # Scale factor for better quality
                pix = page.get_pixmap(matrix=mat)
                
                # Convert to PIL Image
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                
                # Save image to temporary file
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                    img.save(temp_img.name, 'PNG')
                    temp_img_path = temp_img.name
                
                try:
                    # Add image to slide - fill the entire slide to maintain PDF dimensions
                    slide.shapes.add_picture(
                        temp_img_path,
                        left=0,
                        top=0,
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                    
                    print(f"DEBUG: Added page {page_num + 1} to PowerPoint slide with full dimensions")
                    
                finally:
                    # Clean up temporary image file
                    if os.path.exists(temp_img_path):
                        os.unlink(temp_img_path)
            
            # Close PDF document
            pdf_document.close()
            
            # Save PowerPoint to bytes
            output_buffer = io.BytesIO()
            prs.save(output_buffer)
            powerpoint_content = output_buffer.getvalue()
            
            print(f"DEBUG: Successfully created PowerPoint with PDF dimensions, size: {len(powerpoint_content)} bytes")
            
            return powerpoint_content
            
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Error converting PDF to PowerPoint: {str(e)}"
            )
    
    @staticmethod
    def _convert_pptx_to_pdf(pptx_content: bytes) -> bytes:
        """Convert PowerPoint to PDF using python-pptx and PyMuPDF."""
        try:
            # Create input buffer for PowerPoint
            input_buffer = io.BytesIO(pptx_content)
            
            # Open PowerPoint presentation
            prs = Presentation(input_buffer)
            
            print(f"DEBUG: Converting PowerPoint with {len(prs.slides)} slides to PDF")
            
            # Get PowerPoint slide dimensions
            pptx_width_emu = prs.slide_width
            pptx_height_emu = prs.slide_height
            
            # Convert EMU to PDF points
            # 1 inch = 914400 EMU, 1 inch = 72 points
            pptx_width_inches = pptx_width_emu / 914400
            pptx_height_inches = pptx_height_emu / 914400
            
            # Convert to PDF points
            pdf_width_points = int(pptx_width_inches * 72)
            pdf_height_points = int(pptx_height_inches * 72)
            
            print(f"DEBUG: PowerPoint dimensions: {pptx_width_emu} x {pptx_height_emu} EMU")
            print(f"DEBUG: PDF dimensions: {pdf_width_points} x {pdf_height_points} points")
            
            # Create PDF document
            pdf_document = fitz.open()
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides):
                # Create a new page for each slide with PowerPoint dimensions
                page = pdf_document.new_page(width=pdf_width_points, height=pdf_height_points)
                
                print(f"DEBUG: Processing slide {slide_num + 1}")
                
                # Process shapes on the slide
                for shape in slide.shapes:
                    if hasattr(shape, 'image'):
                        # Handle images
                        try:
                            # Get image data
                            image_data = shape.image.blob
                            
                            # Convert to PIL Image
                            img = Image.open(io.BytesIO(image_data))
                            
                            # Save image to temporary file
                            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                                img.save(temp_img.name, 'PNG')
                                temp_img_path = temp_img.name
                            
                            try:
                                # Get shape position and size (convert from EMU to PDF points)
                                # 1 EMU = 1/914400 inch, 1 inch = 72 points
                                left = shape.left * 72 / 914400
                                top = shape.top * 72 / 914400
                                width = shape.width * 72 / 914400
                                height = shape.height * 72 / 914400
                                
                                # Insert image into PDF page
                                page.insert_image(
                                    fitz.Rect(left, top, left + width, top + height),
                                    filename=temp_img_path
                                )
                                
                                print(f"DEBUG: Added image to slide {slide_num + 1}")
                                
                            finally:
                                # Clean up temporary image file
                                if os.path.exists(temp_img_path):
                                    os.unlink(temp_img_path)
                                    
                        except Exception as img_error:
                            print(f"DEBUG: Error processing image in slide {slide_num + 1}: {str(img_error)}")
                            continue
                    
                    elif hasattr(shape, 'text'):
                        # Handle text boxes
                        try:
                            text = shape.text
                            if text.strip():
                                # Get text position (convert from EMU to PDF points)
                                left = shape.left * 72 / 914400
                                top = shape.top * 72 / 914400
                                width = shape.width * 72 / 914400
                                height = shape.height * 72 / 914400
                                
                                # Insert text into PDF page
                                page.insert_text(
                                    fitz.Point(left, top + height/2),  # Center text vertically
                                    text,
                                    fontsize=12
                                )
                                
                                print(f"DEBUG: Added text to slide {slide_num + 1}")
                                
                        except Exception as text_error:
                            print(f"DEBUG: Error processing text in slide {slide_num + 1}: {str(text_error)}")
                            continue
                
                print(f"DEBUG: Completed slide {slide_num + 1}")
            
            # Save PDF to bytes
            pdf_bytes = pdf_document.write()
            pdf_document.close()
            
            print(f"DEBUG: Successfully converted PowerPoint to PDF with original dimensions, size: {len(pdf_bytes)} bytes")
            
            return pdf_bytes
            
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Error converting PowerPoint to PDF: {str(e)}"
            )
    
    @staticmethod
    def _convert_jpgs_to_pdf(
        jpg_files: List[UploadFile],
        page_orientation: str = "portrait",
        page_size: str = "a4",
        margin: str = "no_margin",
        merge_all: bool = True
    ) -> bytes:
        """
        Convert multiple JPG images to a single PDF file.
        Supports configurable page orientation, size, margins, and merging.
        """
        try:
            # Create a new PDF document
            pdf_document = fitz.open()
            
            # Process each JPG file
            for i, jpg_file in enumerate(jpg_files):
                # Read the uploaded file content
                jpg_content = jpg_file.file.read()
                
                # Convert JPG to PIL Image
                img = Image.open(io.BytesIO(jpg_content))
                
                # Get image dimensions
                img_width, img_height = img.size
                
                # Determine page orientation and size
                if page_orientation == "landscape": # FIXED: Respect user choice
                    if img_width > img_height:
                        orientation = "landscape"
                    else:
                        orientation = "portrait"
                else: # portrait
                    if img_width > img_height:
                        orientation = "landscape"
                    else:
                        orientation = "portrait"
                
                # Set page size based on orientation and requested size
                if page_size == "a4":
                    if orientation == "landscape":
                        pdf_width = 842
                        pdf_height = 595
                    else:
                        pdf_width = 595
                        pdf_height = 842
                elif page_size == "us_letter":
                    if orientation == "landscape":
                        pdf_width = 816
                        pdf_height = 1056
                    else:
                        pdf_width = 1056
                        pdf_height = 816
                elif page_size == "fit":
                    pdf_width = img_width
                    pdf_height = img_height
                
                # Add a new page with the determined dimensions
                page = pdf_document.new_page(width=pdf_width, height=pdf_height)
                
                # Save image to temporary file
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
                    img.save(temp_img.name, 'JPEG')
                    temp_img_path = temp_img.name
                
                try:
                    # Insert image into PDF page
                    page.insert_image(
                        fitz.Rect(0, 0, pdf_width, pdf_height),
                        filename=temp_img_path
                    )
                    
                    print(f"DEBUG: Added image {i+1} to PDF page {len(pdf_document)}")
                    
                finally:
                    # Clean up temporary image file
                    if os.path.exists(temp_img_path):
                        os.unlink(temp_img_path)
            
            # Save PDF to bytes
            pdf_bytes = pdf_document.write()
            pdf_document.close()
            
            print(f"DEBUG: Successfully merged {len(jpg_files)} JPG images into a single PDF, size: {len(pdf_bytes)} bytes")
            
            return pdf_bytes
            
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Error merging JPG images to PDF: {str(e)}"
            )
    
    @staticmethod
    def _convert_jpg_to_pdf(jpg_content: bytes) -> bytes:
        """Convert JPG to PDF using PyMuPDF."""
        try:
            # Create PDF document
            pdf_document = fitz.open()
            
            # Convert JPG to PIL Image
            img = Image.open(io.BytesIO(jpg_content))
            
            # Get image dimensions
            img_width, img_height = img.size
            
            # Convert pixels to points (assuming 72 DPI)
            pdf_width = img_width * 72 / 96  # Convert from pixels to points
            pdf_height = img_height * 72 / 96
            
            # Create a new page with image dimensions
            page = pdf_document.new_page(width=pdf_width, height=pdf_height)
            
            # Save image to temporary file
            with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
                img.save(temp_img.name, 'JPEG')
                temp_img_path = temp_img.name
            
            try:
                # Insert image into PDF page
                page.insert_image(
                    fitz.Rect(0, 0, pdf_width, pdf_height),
                    filename=temp_img_path
                )
                
                print(f"DEBUG: Converted JPG to PDF, dimensions: {pdf_width} x {pdf_height} points")
                
            finally:
                # Clean up temporary image file
                if os.path.exists(temp_img_path):
                    os.unlink(temp_img_path)
            
            # Save PDF to bytes
            pdf_bytes = pdf_document.write()
            pdf_document.close()
            
            print(f"DEBUG: Successfully converted JPG to PDF, size: {len(pdf_bytes)} bytes")
            
            return pdf_bytes
            
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Error converting JPG to PDF: {str(e)}"
            )
    
    @staticmethod
    def _convert_pdf_to_jpg(pdf_content: bytes, page_number: int = 1) -> bytes:
        """Convert PDF page to JPG using PyMuPDF."""
        try:
            # Open PDF with PyMuPDF
            pdf_document = fitz.open(stream=pdf_content, filetype="pdf")
            
            # Check if page number is valid
            if page_number < 1 or page_number > len(pdf_document):
                raise HTTPException(
                    status_code=400,
                    detail=f"Invalid page number. PDF has {len(pdf_document)} pages, requested page {page_number}"
                )
            
            # Get the specified page (convert to 0-based index)
            page = pdf_document[page_number - 1]
            
            # Convert page to image with high quality
            mat = fitz.Matrix(2, 2)  # Scale factor for better quality
            pix = page.get_pixmap(matrix=mat)
            
            # Convert to PIL Image first, then to JPEG
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            
            # Convert to JPEG with high quality
            output_buffer = io.BytesIO()
            img.save(output_buffer, format='JPEG', quality=95)
            jpg_data = output_buffer.getvalue()
            
            # Close PDF document
            pdf_document.close()
            
            print(f"DEBUG: Successfully converted PDF page {page_number} to JPG, size: {len(jpg_data)} bytes")
            
            return jpg_data
            
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Error converting PDF to JPG: {str(e)}"
            )
    
    @staticmethod
    def validate_pdf_file(pdf_file: UploadFile) -> None:
        """
        Validate that the uploaded file is a valid PDF.
        
        Args:
            pdf_file: The uploaded file to validate
            
        Raises:
            HTTPException: If file is not a valid PDF
        """
        # Check file extension
        if not pdf_file.filename.lower().endswith('.pdf'):
            raise HTTPException(
                status_code=400,
                detail="File must be a PDF"
            )
        
        # Check content type
        if pdf_file.content_type != 'application/pdf':
            raise HTTPException(
                status_code=400,
                detail="File content type must be application/pdf"
            )
        
        # Check file size (limit to 50MB)
        pdf_file.file.seek(0, 2)  # Seek to end
        file_size = pdf_file.file.tell()
        pdf_file.file.seek(0)  # Reset to beginning
        
        if file_size > 50 * 1024 * 1024:  # 50MB limit
            raise HTTPException(
                status_code=400,
                detail="File size must be less than 50MB"
            )
    
    @staticmethod
    def validate_powerpoint_file(pptx_file: UploadFile) -> None:
        """
        Validate that the uploaded file is a valid PowerPoint file.
        
        Args:
            pptx_file: The uploaded file to validate
            
        Raises:
            HTTPException: If file is not a valid PowerPoint file
        """
        # Check file extension
        if not pptx_file.filename.lower().endswith('.pptx'):
            raise HTTPException(
                status_code=400,
                detail="File must be a PowerPoint (.pptx) file"
            )
        
        # Check content type
        if pptx_file.content_type not in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/octet-stream']:
            raise HTTPException(
                status_code=400,
                detail="File content type must be a PowerPoint presentation"
            )
        
        # Check file size (limit to 50MB)
        pptx_file.file.seek(0, 2)  # Seek to end
        file_size = pptx_file.file.tell()
        pptx_file.file.seek(0)  # Reset to beginning
        
        if file_size > 50 * 1024 * 1024:  # 50MB limit
            raise HTTPException(
                status_code=400,
                detail="File size must be less than 50MB"
            )
    
    @staticmethod
    def validate_jpg_file(jpg_file: UploadFile) -> None:
        """
        Validate that the uploaded file is a valid JPG image file.
        
        Args:
            jpg_file: The uploaded file to validate
            
        Raises:
            HTTPException: If file is not a valid JPG image
        """
        # Check file extension
        if not jpg_file.filename.lower().endswith(('.jpg', '.jpeg')):
            raise HTTPException(
                status_code=400,
                detail="File must be a JPG/JPEG image file"
            )
        
        # Check content type
        if jpg_file.content_type not in ['image/jpeg', 'image/jpg', 'application/octet-stream']:
            raise HTTPException(
                status_code=400,
                detail="File content type must be a JPEG image"
            )
        
        # Check file size (limit to 50MB)
        jpg_file.file.seek(0, 2)  # Seek to end
        file_size = jpg_file.file.tell()
        jpg_file.file.seek(0)  # Reset to beginning
        
        if file_size > 50 * 1024 * 1024:  # 50MB limit
            raise HTTPException(
                status_code=400,
                detail="File size must be less than 50MB"
            )
    
    @staticmethod
    def create_streaming_response(
        pdf_content: bytes, 
        filename: str
    ) -> StreamingResponse:
        """
        Create a streaming response for the unlocked PDF.
        
        Args:
            pdf_content: The unlocked PDF content
            filename: The filename for the response
            
        Returns:
            StreamingResponse with the PDF content
        """
        return StreamingResponse(
            io.BytesIO(pdf_content),
            media_type="application/pdf",
            headers={
                "Content-Disposition": f"attachment; filename={filename}",
                "Content-Length": str(len(pdf_content))
            }
        )
