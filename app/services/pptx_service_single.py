import os
import re
import uuid
import fitz  # PyMuPDF
import requests
from datetime import datetime
from typing import List, Optional, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
from app.models.pdf_model import ConversionRequest, ConversionResponse

class PowerPointServiceSingle:
    """PowerPoint service that creates ONE slide replicating the PDF layout"""
    
    def __init__(self, output_dir: str = "outputs/pptx"):
        # Use /tmp in Lambda environment
        if os.environ.get("LAMBDA_EXECUTION"):
            self.output_dir = "/tmp/outputs/pptx"
        else:
            self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)
    
    def convert_pdf_to_pptx(self, conversion_request: ConversionRequest) -> ConversionResponse:
        """Convert PDF to PowerPoint with SINGLE slide per page"""
        temp_pdf_path = None
        try:
            # Generate unique filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            unique_id = str(uuid.uuid4())[:8]
            safe_name = re.sub(r'[^\w\s-]', '', conversion_request.output_name).strip()
            safe_name = re.sub(r'[-\s]+', '_', safe_name)
            
            filename = f"{safe_name}_{timestamp}_{unique_id}.pptx"
            filepath = os.path.join(self.output_dir, filename)
            
            # Handle URL or local file
            pdf_path = self._get_pdf_path(conversion_request.pdf_path)
            if pdf_path != conversion_request.pdf_path:
                temp_pdf_path = pdf_path  # Mark for cleanup
            
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Process each page of the PDF
            pdf_document = fitz.open(pdf_path)
            
            for page_num in range(len(pdf_document)):
                self._create_single_page_slide(prs, pdf_document, page_num, pdf_path)
            
            pdf_document.close()
            
            # Save presentation
            prs.save(filepath)
            
            # Create HTTP URL for download
            if os.environ.get("LAMBDA_EXECUTION"):
                # Lambda URL will be set by API Gateway
                base_url = os.environ.get("API_BASE_URL", "https://your-api-id.execute-api.region.amazonaws.com/prod")
                http_url = f"{base_url}/convert/download-pptx/{filename}"
            else:
                http_url = f"http://localhost:8000/convert/download-pptx/{filename}"
            
            return ConversionResponse(
                success=True,
                message=f"PowerPoint created successfully: {filename}",
                pptx_path=http_url
            )
            
        except Exception as e:
            return ConversionResponse(
                success=False,
                message=f"Error converting PDF to PowerPoint: {str(e)}"
            )
        finally:
            # Clean up temporary PDF file if downloaded
            if temp_pdf_path and os.path.exists(temp_pdf_path):
                try:
                    os.remove(temp_pdf_path)
                except:
                    pass
    
    def _get_pdf_path(self, pdf_input: str) -> str:
        """Get PDF path - download if URL, return path if local file"""
        # Check if it's a URL
        if pdf_input.startswith(('http://', 'https://')):
            return self._download_pdf_from_url(pdf_input)
        else:
            # It's a local file path
            if not os.path.exists(pdf_input):
                raise FileNotFoundError(f"PDF file not found: {pdf_input}")
            return pdf_input
    
    def _download_pdf_from_url(self, url: str) -> str:
        """Download PDF from URL and return local path"""
        try:
            print(f"📥 Downloading PDF from URL: {url}")
            
            # Create temp directory
            temp_dir = os.path.join(self.output_dir, "temp_pdfs")
            os.makedirs(temp_dir, exist_ok=True)
            
            # Generate temp filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            temp_filename = f"temp_pdf_{timestamp}.pdf"
            temp_path = os.path.join(temp_dir, temp_filename)
            
            # Download with proper headers
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            
            response = requests.get(url, headers=headers, timeout=30, stream=True)
            response.raise_for_status()
            
            # Check if response is actually a PDF
            content_type = response.headers.get('content-type', '').lower()
            if 'pdf' not in content_type and not url.lower().endswith('.pdf'):
                # Try to detect PDF by content
                first_chunk = next(iter(response.iter_content(chunk_size=1024)), b'')
                if not first_chunk.startswith(b'%PDF'):
                    raise ValueError("URL does not point to a valid PDF file")
            
            # Save the PDF
            total_size = 0
            with open(temp_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    if chunk:
                        f.write(chunk)
                        total_size += len(chunk)
            
            print(f"✅ PDF downloaded successfully: {total_size} bytes")
            
            # Validate the downloaded file is a valid PDF
            try:
                test_doc = fitz.open(temp_path)
                page_count = len(test_doc)
                test_doc.close()
                print(f"✅ PDF validation successful: {page_count} page(s)")
            except Exception as e:
                os.remove(temp_path)
                raise ValueError(f"Downloaded file is not a valid PDF: {str(e)}")
            
            return temp_path
            
        except requests.exceptions.RequestException as e:
            raise Exception(f"Failed to download PDF from URL: {str(e)}")
        except Exception as e:
            raise Exception(f"Error processing PDF URL: {str(e)}")
    
    def _create_single_page_slide(self, prs: Presentation, pdf_document, page_num: int, pdf_path: str):
        """Create ONE slide that replicates the entire PDF page"""
        page = pdf_document.load_page(page_num)
        
        # Use blank layout for maximum control
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Get page as high-resolution image
        mat = fitz.Matrix(3.0, 3.0)  # 3x zoom for high quality
        pix = page.get_pixmap(matrix=mat)
        
        # Save page as image
        temp_dir = os.path.join(self.output_dir, "temp_page_images")
        os.makedirs(temp_dir, exist_ok=True)
        
        page_img_path = os.path.join(temp_dir, f"page_{page_num + 1}.png")
        pix.save(page_img_path)
        
        try:
            # Add the full page image to the slide
            # Calculate size to fit slide while maintaining aspect ratio
            slide_width = 10.0  # PowerPoint slide width in inches
            slide_height = 7.5   # PowerPoint slide height in inches
            
            # Get image dimensions
            img_width_px = pix.width
            img_height_px = pix.height
            aspect_ratio = img_width_px / img_height_px
            
            # Calculate best fit
            if aspect_ratio > slide_width / slide_height:
                # Image is wider - fit to width
                width = Inches(slide_width * 0.95)  # 95% to leave small margins
                height = Inches((slide_width * 0.95) / aspect_ratio)
            else:
                # Image is taller - fit to height
                height = Inches(slide_height * 0.95)
                width = Inches((slide_height * 0.95) * aspect_ratio)
            
            # Center the image on the slide
            left = Inches((slide_width - width.inches) / 2)
            top = Inches((slide_height - height.inches) / 2)
            
            # Add the page image to slide
            slide.shapes.add_picture(page_img_path, left, top, width, height)
            
            # Add minimal title at the top if desired
            if page_num == 0:  # Only for first page
                title_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(9.6), Inches(0.4))
                title_frame = title_box.text_frame
                title_frame.text = f"Document: {os.path.basename(pdf_path)}"
                
                # Style the title
                for paragraph in title_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
                        run.font.color.rgb = RGBColor(64, 64, 64)  # Dark gray
            
        except Exception as e:
            # Fallback: create error slide
            error_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
            error_frame = error_box.text_frame
            error_frame.text = f"Error loading page {page_num + 1}: {str(e)}"
            
        finally:
            # Clean up
            pix = None
            if os.path.exists(page_img_path):
                os.remove(page_img_path)
    
    def create_sample_pptx(self, title: str, content: List[str]) -> str:
        """Create a sample PowerPoint presentation"""
        try:
            # Generate unique filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            unique_id = str(uuid.uuid4())[:8]
            safe_title = re.sub(r'[^\w\s-]', '', title).strip()
            safe_title = re.sub(r'[-\s]+', '_', safe_title)
            
            filename = f"{safe_title}_{timestamp}_{unique_id}.pptx"
            filepath = os.path.join(self.output_dir, filename)
            
            # Create presentation
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            slide.shapes.placeholders[1].text = "Single-Slide PDF Replication Method"
            
            # Content slides
            bullet_slide_layout = prs.slide_layouts[1]
            for i, item in enumerate(content):
                slide = prs.slides.add_slide(bullet_slide_layout)
                slide.shapes.title.text = f"Content Item {i+1}"
                
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.text = item
            
            # Save presentation
            prs.save(filepath)
            
            # Return absolute path
            return os.path.abspath(filepath)
            
        except Exception as e:
            raise Exception(f"Error creating sample PowerPoint: {str(e)}")
