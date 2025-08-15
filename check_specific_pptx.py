#!/usr/bin/env python3
"""
Check specific PowerPoint content
"""

from pptx import Presentation
import sys

def check_pptx(pptx_path):
    """Check the content of a specific PowerPoint file"""
    try:
        prs = Presentation(pptx_path)
        print(f"📊 PowerPoint: {pptx_path}")
        print(f"Total slides: {len(prs.slides)}")
        print("=" * 60)
        
        for i, slide in enumerate(prs.slides):
            print(f"\n📄 Slide {i+1}:")
            
            # Check title
            if slide.shapes.title:
                title = slide.shapes.title.text
                print(f"   Title: {title}")
            
            # Check content
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text and text != title:
                            print(f"   Content: {text}")
            
            print()
        
    except Exception as e:
        print(f"❌ Error: {e}")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
    else:
        pptx_path = "outputs/pptx/improved_conversion_20250815_135827_936920_cedfe053.pptx"
    
    check_pptx(pptx_path)
