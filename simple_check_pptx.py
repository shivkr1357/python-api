#!/usr/bin/env python3
"""
Simple script to check PowerPoint content
"""

import sys
from pptx import Presentation

def check_pptx(pptx_path):
    """Check PowerPoint content"""
    try:
        prs = Presentation(pptx_path)
        
        print(f"📊 PowerPoint: {pptx_path}")
        print(f"Total slides: {len(prs.slides)}")
        print("=" * 60)
        
        for i, slide in enumerate(prs.slides):
            print(f"\n📄 Slide {i+1}:")
            
            # Count shapes
            shape_count = len(slide.shapes)
            print(f"   Total shapes: {shape_count}")
            
            # Check for images
            image_count = 0
            text_count = 0
            
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    image_count += 1
                elif hasattr(shape, 'text_frame') and shape.text_frame:
                    text_count += 1
                elif str(type(shape)).find('Picture') != -1:
                    image_count += 1
            
            print(f"   Images: {image_count}")
            print(f"   Text elements: {text_count}")
            
        return True
        
    except Exception as e:
        print(f"❌ Error: {e}")
        return False

if __name__ == "__main__":
    if len(sys.argv) > 1:
        check_pptx(sys.argv[1])
    else:
        print("Usage: python simple_check_pptx.py <pptx_file>")
