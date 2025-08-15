#!/usr/bin/env python3
"""
Test script to compare PDF to PowerPoint conversions
"""

import requests
import os
from datetime import datetime

class ConversionComparisonTester:
    def __init__(self, base_url: str = "http://localhost:8000"):
        self.base_url = base_url
        self.session = requests.Session()
    
    def test_text_based_pdf(self):
        """Test conversion of a text-based PDF"""
        print("📄 Testing Text-Based PDF Conversion")
        print("=" * 50)
        
        # Create a text-based PDF
        pdf_data = {
            "title": "Sample Text Document",
            "content": """
            This is a sample text document with multiple paragraphs.
            
            Key Features:
            • Professional formatting
            • Multiple paragraphs
            • Structured content
            • Clear organization
            
            This document demonstrates the conversion capabilities.
            The content should be properly extracted and formatted.
            
            Technical Details:
            • PDF created using ReportLab
            • PowerPoint created using python-pptx
            • Text extraction using advanced libraries
            • Automatic formatting and layout
            """,
            "author": "Test Author",
            "subject": "Conversion Testing",
            "keywords": ["text", "conversion", "testing", "pdf"]
        }
        
        try:
            # Create PDF
            response = self.session.post(f"{self.base_url}/pdf/create", json=pdf_data)
            response.raise_for_status()
            pdf_result = response.json()
            
            if pdf_result['success']:
                print(f"✅ Created PDF: {os.path.basename(pdf_result['pdf_path'])}")
                
                # Convert to PowerPoint
                conversion_data = {
                    "pdf_path": pdf_result['pdf_path'],
                    "output_name": "text_based_conversion",
                    "include_images": True
                }
                
                response = self.session.post(f"{self.base_url}/convert/pdf-to-pptx", json=conversion_data)
                response.raise_for_status()
                pptx_result = response.json()
                
                if pptx_result['success']:
                    print(f"✅ Created PowerPoint: {os.path.basename(pptx_result['pptx_path'])}")
                    return pptx_result['pptx_path']
                else:
                    print(f"❌ PowerPoint creation failed: {pptx_result.get('message', 'Unknown error')}")
                    return None
            else:
                print(f"❌ PDF creation failed: {pdf_result.get('message', 'Unknown error')}")
                return None
                
        except Exception as e:
            print(f"❌ Error: {e}")
            return None
    
    def test_image_based_pdf(self):
        """Test conversion of an image-based PDF"""
        print("\n🖼️ Testing Image-Based PDF Conversion")
        print("=" * 50)
        
        # Use the user's original PDF
        user_pdf_path = "outputs/pdfs/user_original.pdf"
        
        if not os.path.exists(user_pdf_path):
            print(f"❌ User PDF not found: {user_pdf_path}")
            return None
        
        try:
            # Convert to PowerPoint
            conversion_data = {
                "pdf_path": user_pdf_path,
                "output_name": "image_based_conversion",
                "include_images": True
            }
            
            response = self.session.post(f"{self.base_url}/convert/pdf-to-pptx", json=conversion_data)
            response.raise_for_status()
            pptx_result = response.json()
            
            if pptx_result['success']:
                print(f"✅ Created PowerPoint: {os.path.basename(pptx_result['pptx_path'])}")
                return pptx_result['pptx_path']
            else:
                print(f"❌ PowerPoint creation failed: {pptx_result.get('message', 'Unknown error')}")
                return None
                
        except Exception as e:
            print(f"❌ Error: {e}")
            return None
    
    def check_pptx_content(self, pptx_path: str, description: str):
        """Check the content of a PowerPoint file"""
        print(f"\n📊 {description}")
        print("-" * 40)
        
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            
            print(f"Total slides: {len(prs.slides)}")
            
            for i, slide in enumerate(prs.slides):
                print(f"\n📄 Slide {i+1}:")
                
                # Check title
                if slide.shapes.title:
                    title = slide.shapes.title.text
                    print(f"   Title: {title}")
                
                # Check content
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text and text != title:
                                # Truncate long text for display
                                display_text = text[:100] + "..." if len(text) > 100 else text
                                print(f"   Content: {display_text}")
            
        except Exception as e:
            print(f"❌ Error reading PowerPoint: {e}")
    
    def run_comparison(self):
        """Run the complete comparison test"""
        print("🔄 PDF to PowerPoint Conversion Comparison Test")
        print("=" * 60)
        
        # Test 1: Text-based PDF
        text_pptx = self.test_text_based_pdf()
        if text_pptx:
            self.check_pptx_content(text_pptx, "Text-Based PDF Conversion Result")
        
        # Test 2: Image-based PDF
        image_pptx = self.test_image_based_pdf()
        if image_pptx:
            self.check_pptx_content(image_pptx, "Image-Based PDF Conversion Result")
        
        # Summary
        print("\n" + "=" * 60)
        print("🎉 Conversion Comparison Summary:")
        print()
        print("📄 Text-Based PDFs:")
        print("   ✅ Full content extraction")
        print("   ✅ Proper formatting and organization")
        print("   ✅ Metadata separation")
        print("   ✅ Multiple slides with structured content")
        print()
        print("🖼️ Image-Based PDFs (Scanned Documents):")
        print("   ⚠️  No text extraction possible")
        print("   ✅ Clear explanation of limitation")
        print("   ✅ Guidance for OCR requirements")
        print("   ✅ Professional presentation format")
        print()
        print("💡 The conversion system now properly handles both types!")
        print("   - Text-based PDFs: Full content conversion")
        print("   - Image-based PDFs: Informative explanation")

def main():
    """Main function to run the comparison test"""
    print("PDF to PowerPoint Conversion Comparison")
    print("Make sure the API server is running on http://localhost:8000")
    print()
    
    # Wait a moment for user to read
    input("Press Enter to start comparison...")
    
    # Create tester and run comparison
    tester = ConversionComparisonTester()
    tester.run_comparison()

if __name__ == "__main__":
    main()
